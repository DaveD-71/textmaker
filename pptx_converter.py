from __future__ import annotations

import argparse
import json
import logging
import re
import shutil
from collections import defaultdict
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.dml.color import ColorFormat
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

EMU_PER_PIXEL = 9525


LOGGER = logging.getLogger("pptx_converter")


def sanitize_filename(name: str) -> str:
    """Convert presentation title to safe folder name."""
    name = re.sub(r"[^\w\s-]", "", name)
    name = re.sub(r"[-\s]+", "_", name)
    sanitized = name.strip("_").lower()
    return sanitized or "presentation"


def emu_to_px(value: int | None) -> int:
    if value is None:
        return 0
    return int(value) // EMU_PER_PIXEL


def safe_get(obj: Any, attr: str, default: Any = None) -> Any:
    """Safely get attribute from object, return default if None or error."""
    try:
        value = getattr(obj, attr, default)
        return value if value is not None else default
    except Exception:
        return default


def rgb_to_hex(rgb_obj: Any) -> str | None:
    """Convert RGBColor object to hex string #RRGGBB."""
    if rgb_obj is None:
        return None
    try:
        return f"#{rgb_obj[0]:02X}{rgb_obj[1]:02X}{rgb_obj[2]:02X}"
    except Exception:
        return None


def color_format_to_hex(color: ColorFormat | None) -> str | None:
    if color is None:
        return None
    rgb = safe_get(color, "rgb")
    if rgb is not None:
        return rgb_to_hex(rgb)
    return None


def extract_image(shape: Any, slide_num: int, image_num: int, output_dir: Path) -> str:
    """Extract image from shape and save to output_dir/images/."""
    image = shape.image
    extension = image.ext
    filename = f"slide{slide_num}_image{image_num}.{extension}"
    filepath = output_dir / "images" / filename
    with filepath.open("wb") as fh:
        fh.write(image.blob)
    return f"images/{filename}"


def extract_media_target(shape: Any) -> str | None:
    """Best-effort extraction of media relationship target name from shape XML."""
    try:
        blips = shape.element.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}videoFile")
        if blips:
            embed = blips[0].attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link")
            if embed:
                rel = shape.part.related_part(embed)
                return rel.partname
    except Exception:
        return None
    return None


def get_alt_text(shape: Any) -> str:
    descr = ""
    try:
        c_nv_pr = shape.element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
        if c_nv_pr is not None:
            descr = c_nv_pr.attrib.get("descr", "")
    except Exception:
        descr = ""
    if descr:
        return descr
    return safe_get(shape, "name", "") or ""


def detect_text_direction(shape: Any) -> str:
    try:
        if safe_get(shape, "rotation", 0) in (90, 270):
            return "vertical"
        body_pr = shape.text_frame._txBody.bodyPr  # pylint: disable=protected-access
        vert = body_pr.attrib.get("vert") if body_pr is not None else None
        return "vertical" if vert and vert != "horz" else "horizontal"
    except Exception:
        return "horizontal"


def map_alignment(alignment: PP_ALIGN | None) -> str:
    mapping = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justified",
        PP_ALIGN.JUSTIFY_LOW: "justified",
        PP_ALIGN.DISTRIBUTE: "justified",
        PP_ALIGN.THAI_DISTRIBUTE: "justified",
    }
    return mapping.get(alignment, "left")


def map_vertical_anchor(anchor: MSO_ANCHOR | None) -> str:
    mapping = {
        MSO_ANCHOR.TOP: "top",
        MSO_ANCHOR.MIDDLE: "middle",
        MSO_ANCHOR.BOTTOM: "bottom",
    }
    return mapping.get(anchor, "top")


def extract_fill(shape: Any) -> dict[str, Any] | None:
    fill = safe_get(shape, "fill")
    if fill is None:
        return None
    fill_type = safe_get(fill, "type")
    if fill_type == MSO_FILL_TYPE.SOLID:
        return {"type": "solid", "color": color_format_to_hex(safe_get(fill, "fore_color"))}
    if fill_type == MSO_FILL_TYPE.GRADIENT:
        stops = []
        try:
            for stop in fill.gradient_stops:
                stops.append({
                    "position": safe_get(stop, "position", 0),
                    "color": color_format_to_hex(safe_get(stop, "color")),
                })
        except Exception:
            stops = []
        return {"type": "gradient", "stops": stops}
    if fill_type == MSO_FILL_TYPE.PATTERNED:
        return {"type": "pattern"}
    return None


def extract_border(shape: Any) -> dict[str, Any] | None:
    line = safe_get(shape, "line")
    if line is None:
        return None
    color = None
    try:
        color = color_format_to_hex(line.color)
    except Exception:
        color = None
    width = safe_get(safe_get(line, "width"), "pt", 1)
    return {"color": color, "width": width or 1, "style": "solid"}


def extract_background(slide: Any) -> dict[str, Any] | None:
    fill = safe_get(safe_get(slide, "background"), "fill")
    if fill is None:
        return None
    fill_type = safe_get(fill, "type")
    if fill_type == MSO_FILL_TYPE.SOLID:
        return {"type": "solid", "color": color_format_to_hex(safe_get(fill, "fore_color"))}
    if fill_type == MSO_FILL_TYPE.GRADIENT:
        return {"type": "gradient"}
    if fill_type == MSO_FILL_TYPE.PATTERNED:
        return {"type": "pattern"}
    if fill_type == MSO_FILL_TYPE.PICTURE:
        return {"type": "picture"}
    return None


def extract_hyperlink(shape: Any) -> str | None:
    try:
        link = shape.click_action.hyperlink
        return safe_get(link, "address")
    except Exception:
        return None


def extract_bullet(paragraph: Any) -> dict[str, Any] | None:
    try:
        level = safe_get(paragraph, "level", 0)
        ppr = paragraph._p.pPr  # pylint: disable=protected-access
        if ppr is None:
            return None
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        if ppr.find(f"{ns}buNone") is not None:
            return None
        char_el = ppr.find(f"{ns}buChar")
        if char_el is not None:
            return {"enabled": True, "level": level, "character": char_el.attrib.get("char", "•")}
        if level and level >= 0:
            return {"enabled": True, "level": level, "character": "•"}
    except Exception:
        return None
    return None


def extract_chart_data(chart: Any) -> dict[str, Any]:
    categories: list[str] = []
    series_data: list[dict[str, Any]] = []
    try:
        if chart.plots and chart.plots[0].categories:
            categories = [str(c.label) for c in chart.plots[0].categories]
    except Exception:
        categories = []

    try:
        for series in chart.series:
            values = []
            try:
                values = [v for v in series.values]
            except Exception:
                values = []
            series_data.append({"name": safe_get(series, "name"), "values": values})
    except Exception:
        pass
    return {"categories": categories, "series": series_data}


def parse_slide_xml_for_animations(slide: Any) -> dict[str, Any]:
    # MVP: return null/empty when we cannot confidently parse all details.
    _ = slide
    return {"transition": None, "animations": {}}


def create_output_structure(base_output_dir: Path, pres_name: str, overwrite: bool = False) -> Path:
    out_dir = base_output_dir / sanitize_filename(pres_name)
    if out_dir.exists():
        if overwrite:
            shutil.rmtree(out_dir)
        else:
            raise FileExistsError(
                f"Output directory already exists: {out_dir}. Use --overwrite to replace it."
            )
    (out_dir / "images").mkdir(parents=True, exist_ok=True)
    return out_dir


def _shape_kind(shape: Any) -> str:
    if safe_get(shape, "has_text_frame", False):
        return "text_box"
    if safe_get(shape, "shape_type") == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if safe_get(shape, "has_table", False):
        return "table"
    if safe_get(shape, "has_chart", False):
        return "chart"
    if safe_get(shape, "shape_type") == MSO_SHAPE_TYPE.MEDIA:
        return "video"
    return "shape"


def convert_pptx(pptx_path: Path, output_dir: Path, overwrite: bool = False) -> tuple[int, int]:
    prs = Presentation(str(pptx_path))
    title = safe_get(prs.core_properties, "title") or pptx_path.stem
    out_dir = create_output_structure(output_dir, title, overwrite=overwrite)

    content_yaml: dict[str, Any] = {"presentation": {"title": title}, "slides": []}
    slide_properties: dict[str, Any] = {
        "presentation": {
            "slide_width": emu_to_px(prs.slide_width),
            "slide_height": emu_to_px(prs.slide_height),
        },
        "slides": {},
    }
    object_properties: dict[str, Any] = {"slides": {}}
    text_formatting: dict[str, Any] = {"slides": {}}
    shape_styling: dict[str, Any] = {"slides": {}}
    animations: dict[str, Any] = {"slides": {}}
    embedded_objects: dict[str, Any] = {"slides": {}}

    total_images = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        counters = defaultdict(int)
        slide_key = str(slide_idx)

        slide_entry: dict[str, Any] = {
            "slide_number": slide_idx,
            "layout": safe_get(slide.slide_layout, "name", "Unknown"),
            "text_boxes": [],
            "objects": [],
        }

        slide_properties["slides"][slide_key] = {
            "layout_name": safe_get(slide.slide_layout, "name", "Unknown"),
            "background": extract_background(slide),
            "slide_number_visible": False,
            "hidden": False,
        }

        object_properties["slides"][slide_key] = {"text_boxes": {}, "objects": {}}
        text_formatting["slides"][slide_key] = {"text_boxes": {}}
        shape_styling["slides"][slide_key] = {"text_boxes": {}, "objects": {}}
        embedded_objects["slides"][slide_key] = {
            "tables": {},
            "charts": {},
            "videos": {},
            "audio": {},
            "smartart": {},
            "embedded_documents": {},
        }
        animations["slides"][slide_key] = parse_slide_xml_for_animations(slide)

        for z_order, shape in enumerate(slide.shapes, start=1):
            kind = _shape_kind(shape)
            if kind == "text_box":
                counters["tb"] += 1
                obj_id = f"tb_{counters['tb']}"
                text_content = safe_get(shape, "text", "")
                slide_entry["text_boxes"].append({"id": obj_id, "content": text_content or ""})

                object_properties["slides"][slide_key]["text_boxes"][obj_id] = {
                    "name": safe_get(shape, "name"),
                    "position": {
                        "x": emu_to_px(safe_get(shape, "left")),
                        "y": emu_to_px(safe_get(shape, "top")),
                        "width": emu_to_px(safe_get(shape, "width")),
                        "height": emu_to_px(safe_get(shape, "height")),
                    },
                    "z_order": z_order,
                    "rotation": safe_get(shape, "rotation", 0) or 0,
                    "visible": safe_get(shape, "visible", True),
                    "locked": False,
                    "hyperlink": extract_hyperlink(shape),
                }

                paragraph = None
                run = None
                tf = safe_get(shape, "text_frame")
                if tf is not None and tf.paragraphs:
                    paragraph = tf.paragraphs[0]
                    if paragraph.runs:
                        run = paragraph.runs[0]

                font = safe_get(run, "font")
                font_size = safe_get(safe_get(font, "size"), "pt")
                text_formatting["slides"][slide_key]["text_boxes"][obj_id] = {
                    "font_name": safe_get(font, "name", "Calibri") or "Calibri",
                    "font_size": font_size,
                    "font_color": color_format_to_hex(safe_get(font, "color")) or "#000000",
                    "bold": bool(safe_get(font, "bold", False)),
                    "italic": bool(safe_get(font, "italic", False)),
                    "underline": bool(safe_get(font, "underline", False)),
                    "alignment": map_alignment(safe_get(paragraph, "alignment")),
                    "vertical_alignment": map_vertical_anchor(safe_get(tf, "vertical_anchor")),
                    "line_spacing": safe_get(paragraph, "line_spacing", 1.0) if paragraph else 1.0,
                    "paragraph_spacing_before": safe_get(safe_get(paragraph, "space_before"), "pt", 0),
                    "paragraph_spacing_after": safe_get(safe_get(paragraph, "space_after"), "pt", 0),
                    "bullet_style": extract_bullet(paragraph) if paragraph else None,
                    "text_direction": detect_text_direction(shape),
                }

                shape_styling["slides"][slide_key]["text_boxes"][obj_id] = {
                    "fill": extract_fill(shape),
                    "border": extract_border(shape),
                    "shadow": None,
                    "reflection": None,
                    "glow": None,
                    "soft_edges": None,
                    "3d_format": None,
                    "opacity": 1.0,
                }
                continue

            if kind == "image":
                counters["img"] += 1
                obj_id = f"img_{counters['img']}"
                total_images += 1
                source = extract_image(shape, slide_idx, counters["img"], out_dir)
                slide_entry["objects"].append({"id": obj_id, "type": "image", "source": source})
            elif kind == "table":
                counters["table"] += 1
                obj_id = f"table_{counters['table']}"
                table = shape.table
                slide_entry["objects"].append(
                    {
                        "id": obj_id,
                        "type": "table",
                        "rows": len(table.rows),
                        "columns": len(table.columns),
                    }
                )
                table_data = [[cell.text for cell in row.cells] for row in table.rows]
                embedded_objects["slides"][slide_key]["tables"][obj_id] = {
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "data": table_data,
                    "cell_formatting": {},
                }
            elif kind == "chart":
                counters["chart"] += 1
                obj_id = f"chart_{counters['chart']}"
                chart = shape.chart
                chart_type = safe_get(safe_get(chart, "chart_type"), "name")
                slide_entry["objects"].append(
                    {"id": obj_id, "type": "chart", "chart_type": chart_type.lower() if chart_type else None}
                )
                chart_title = None
                try:
                    chart_title = chart.chart_title.text_frame.text if chart.has_title else None
                except Exception:
                    chart_title = None
                embedded_objects["slides"][slide_key]["charts"][obj_id] = {
                    "chart_type": chart_type.lower() if chart_type else None,
                    "title": chart_title,
                    "data": extract_chart_data(chart),
                }
            elif kind == "video":
                counters["video"] += 1
                obj_id = f"video_{counters['video']}"
                media_source = extract_media_target(shape)
                slide_entry["objects"].append(
                    {
                        "id": obj_id,
                        "type": "video",
                        "source": str(media_source) if media_source else None,
                    }
                )
                embedded_objects["slides"][slide_key]["videos"][obj_id] = {
                    "source": str(media_source) if media_source else None,
                    "autoplay": False,
                    "loop": False,
                    "volume": 1.0,
                }
            else:
                counters["shape"] += 1
                obj_id = f"shape_{counters['shape']}"
                slide_entry["objects"].append(
                    {
                        "id": obj_id,
                        "type": str(safe_get(safe_get(shape, "shape_type"), "name", "shape")).lower(),
                    }
                )

            object_properties["slides"][slide_key]["objects"][obj_id] = {
                "name": safe_get(shape, "name"),
                "position": {
                    "x": emu_to_px(safe_get(shape, "left")),
                    "y": emu_to_px(safe_get(shape, "top")),
                    "width": emu_to_px(safe_get(shape, "width")),
                    "height": emu_to_px(safe_get(shape, "height")),
                },
                "z_order": z_order,
                "rotation": safe_get(shape, "rotation", 0) or 0,
                "visible": safe_get(shape, "visible", True),
                "locked": False,
                "hyperlink": extract_hyperlink(shape),
                "alt_text": get_alt_text(shape),
            }

            shape_styling["slides"][slide_key]["objects"][obj_id] = {
                "fill": extract_fill(shape),
                "border": extract_border(shape),
                "shadow": None,
                "reflection": None,
                "glow": None,
                "soft_edges": None,
                "3d_format": None,
                "opacity": 1.0,
            }

        if safe_get(slide, "has_notes_slide", False):
            notes = safe_get(safe_get(slide, "notes_slide"), "notes_text_frame")
            note_text = safe_get(notes, "text") if notes is not None else None
            if note_text:
                slide_entry["notes"] = note_text

        content_yaml["slides"].append(slide_entry)

    files_to_write = {
        "content.yaml": lambda p: yaml.safe_dump(content_yaml, p, sort_keys=False, allow_unicode=True),
        "slide_properties.json": lambda p: json.dump(slide_properties, p, indent=2, ensure_ascii=False),
        "object_properties.json": lambda p: json.dump(object_properties, p, indent=2, ensure_ascii=False),
        "text_formatting.json": lambda p: json.dump(text_formatting, p, indent=2, ensure_ascii=False),
        "shape_styling.json": lambda p: json.dump(shape_styling, p, indent=2, ensure_ascii=False),
        "animations.json": lambda p: json.dump(animations, p, indent=2, ensure_ascii=False),
        "embedded_objects.json": lambda p: json.dump(embedded_objects, p, indent=2, ensure_ascii=False),
    }

    for name, writer in files_to_write.items():
        path = out_dir / name
        with path.open("w", encoding="utf-8") as fh:
            writer(fh)
            fh.write("\n")

    return len(prs.slides), total_images


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Convert PPTX into structured YAML/JSON package.")
    parser.add_argument("input", help="Path to input .pptx file")
    parser.add_argument(
        "--output-dir",
        default=".",
        help="Directory where the output presentation folder will be created (default: current directory)",
    )
    parser.add_argument(
        "--overwrite",
        action="store_true",
        help="Overwrite existing output folder when it already exists",
    )
    parser.add_argument("--verbose", action="store_true", help="Enable verbose logging")
    return parser


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s: %(message)s",
    )

    input_path = Path(args.input)
    if not input_path.exists() or input_path.suffix.lower() != ".pptx":
        parser.error("input must be an existing .pptx file")

    try:
        slide_count, image_count = convert_pptx(
            pptx_path=input_path,
            output_dir=Path(args.output_dir),
            overwrite=args.overwrite,
        )
        print(f"Converted {slide_count} slides, extracted {image_count} images")
    except Exception as exc:
        LOGGER.exception("Conversion failed: %s", exc)
        raise SystemExit(1) from exc


if __name__ == "__main__":
    main()
