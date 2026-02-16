from __future__ import annotations

import json
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx_converter import convert_pptx, emu_to_px, sanitize_filename


_MINI_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc``\xf8\xcf\x00\x00"
    b"\x02\x05\x01\x02\x8d\xc8\x1f\x8f\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _build_sample_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Hello world"

    image_path = path.parent / "tiny.png"
    image_path.write_bytes(_MINI_PNG)
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(1), Inches(1))

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(3), Inches(1.2))
    table_shape.table.cell(0, 0).text = "A"
    table_shape.table.cell(0, 1).text = "B"

    prs.save(path)


def test_utils() -> None:
    assert sanitize_filename("My Presentation 2026!") == "my_presentation_2026"
    assert emu_to_px(9525) == 1


def test_convert_pptx_outputs(tmp_path: Path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    _build_sample_pptx(pptx_path)

    slide_count, image_count = convert_pptx(pptx_path, tmp_path)
    assert slide_count == 1
    assert image_count == 1

    out_dir = tmp_path / "sample"
    assert (out_dir / "images").exists()

    content = yaml.safe_load((out_dir / "content.yaml").read_text(encoding="utf-8"))
    assert content["slides"][0]["text_boxes"][0]["id"] == "tb_1"
    assert content["slides"][0]["objects"][0]["id"] == "img_1"
    assert content["slides"][0]["objects"][1]["id"] == "table_1"

    object_props = json.loads((out_dir / "object_properties.json").read_text(encoding="utf-8"))
    assert "tb_1" in object_props["slides"]["1"]["text_boxes"]
    assert "img_1" in object_props["slides"]["1"]["objects"]
    assert "table_1" in object_props["slides"]["1"]["objects"]

    embedded = json.loads((out_dir / "embedded_objects.json").read_text(encoding="utf-8"))
    assert embedded["slides"]["1"]["tables"]["table_1"]["rows"] == 2
